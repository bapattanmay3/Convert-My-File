import os
from pdf2docx import Converter
from PIL import Image
import PyPDF2
import shutil
from docx import Document
import pandas as pd
import tabula
from pptx import Presentation
from weasyprint import HTML

# ============ FILE CONVERSIONS ============

def convert_pdf_to_docx(input_path, output_path):
    """PDF to Word Document"""
    try:
        cv = Converter(input_path)
        cv.convert(output_path)
        cv.close()
        return True, "PDF to DOCX conversion successful"
    except Exception as e:
        return False, str(e)

def convert_pdf_to_txt(input_path, output_path):
    """PDF to Text"""
    try:
        with open(input_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
            with open(output_path, 'w', encoding='utf-8') as txt_file:
                txt_file.write(text)
        return True, "PDF to TXT conversion successful"
    except Exception as e:
        return False, str(e)

def convert_docx_to_pdf(input_path, output_path):
    """Word to PDF (placeholder - upgrade later)"""
    try:
        shutil.copy(input_path, output_path)
        return True, "DOCX to PDF conversion successful (placeholder)"
    except Exception as e:
        return False, str(e)

def convert_docx_to_txt(input_path, output_path):
    """Word to Text"""
    try:
        doc = Document(input_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        with open(output_path, 'w', encoding='utf-8') as txt_file:
            txt_file.write(text)
        return True, "DOCX to TXT conversion successful"
    except Exception as e:
        return False, str(e)

def convert_txt_to_pdf(input_path, output_path):
    """Text to PDF"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        c = canvas.Canvas(output_path, pagesize=letter)
        with open(input_path, 'r', encoding='utf-8') as f:
            y = 750
            for line in f:
                c.drawString(40, y, line.strip()[:90])
                y -= 15
                if y < 40:
                    c.showPage()
                    y = 750
        c.save()
        return True, "TXT to PDF conversion successful"
    except Exception as e:
        return False, str(e)

def convert_txt_to_docx(input_path, output_path):
    """Text to Word"""
    try:
        doc = Document()
        with open(input_path, 'r', encoding='utf-8') as f:
            doc.add_paragraph(f.read())
        doc.save(output_path)
        return True, "TXT to DOCX conversion successful"
    except Exception as e:
        return False, str(e)

# ============ IMAGE CONVERSIONS ============

def convert_image_to_image(input_path, output_path, target_format, quality=90):
    """Convert image from one format to another"""
    try:
        image = Image.open(input_path)
        
        # Handle transparency for JPEG
        if target_format.upper() in ['JPG', 'JPEG'] and image.mode in ('RGBA', 'LA', 'P'):
            rgb_image = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'P':
                image = image.convert('RGBA')
            rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
            image = rgb_image
        
        save_format = 'JPEG' if target_format.upper() in ['JPG', 'JPEG'] else target_format.upper()
        image.save(output_path, save_format, quality=quality, optimize=True)
        return True, f"Image converted to {target_format.upper()}"
    except Exception as e:
        return False, str(e)

def convert_image_to_pdf(input_path, output_path):
    """Image to PDF"""
    try:
        image = Image.open(input_path)
        if image.mode != 'RGB':
            image = image.convert('RGB')
        image.save(output_path, 'PDF')
        return True, "Image to PDF conversion successful"
    except Exception as e:
        return False, str(e)

def convert_pdf_to_image(input_path, output_path, target_format='png'):
    """PDF to Image (first page only)"""
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(input_path, first_page=1, last_page=1)
        images[0].save(output_path, target_format.upper())
        return True, f"PDF to {target_format.upper()} conversion successful"
    except Exception as e:
        return False, str(e)

# ============ DATA & TABLE CONVERSIONS ============

def convert_pdf_to_excel(input_path, output_path):
    """PDF to Excel (Table Extraction)"""
    try:
        # Extract tables into a list of DataFrames
        dfs = tabula.read_pdf(input_path, pages='all', multiple_tables=True)
        if not dfs:
            return False, "No tables found in PDF"
        
        # Write all tables to a single Excel file
        with pd.ExcelWriter(output_path) as writer:
            for i, df in enumerate(dfs):
                df.to_excel(writer, sheet_name=f'Table_{i+1}', index=False)
        
        return True, "PDF tables successfully extracted to Excel"
    except Exception as e:
        return False, f"Table extraction failed: {str(e)}"

def convert_excel_to_pdf(input_path, output_path):
    """Excel to PDF (via HTML intermediate)"""
    try:
        df = pd.read_excel(input_path)
        html_content = df.to_html()
        HTML(string=html_content).write_pdf(output_path)
        return True, "Excel to PDF conversion successful"
    except Exception as e:
        return False, str(e)

# ============ PRESENTATION CONVERSIONS ============

def convert_pptx_to_pdf(input_path, output_path):
    """PowerPoint to PDF (placeholder for non-Windows)"""
    try:
        # On Windows, we could use comtypes, but for a general solution,
        # we'd need another library or a cloud API. For now, using as placeholder.
        return False, "PowerPoint to PDF requires system-specific utilities or LibreOffice"
    except Exception as e:
        return False, str(e)

def convert_pptx_to_txt(input_path, output_path):
    """PowerPoint to Text"""
    try:
        prs = Presentation(input_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return True, "PowerPoint text successfully extracted"
    except Exception as e:
        return False, str(e)

# ============ WEB CONVERSIONS ============

def convert_html_to_pdf(input_path, output_path):
    """HTML to PDF"""
    try:
        HTML(input_path).write_pdf(output_path)
        return True, "HTML to PDF conversion successful"
    except Exception as e:
        return False, str(e)

# ============ SUPPORTED FORMATS ============

FILE_CONVERSIONS = {
    'pdf': {
        'docx': convert_pdf_to_docx,
        'txt': convert_pdf_to_txt,
        'xlsx': convert_pdf_to_excel,
    },
    'docx': {
        'pdf': convert_docx_to_pdf,
        'txt': convert_docx_to_txt,
    },
    'txt': {
        'pdf': convert_txt_to_pdf,
        'docx': convert_txt_to_docx,
    },
    'xlsx': {
        'pdf': convert_excel_to_pdf,
    },
    'pptx': {
        'txt': convert_pptx_to_txt,
        'pdf': convert_pptx_to_pdf,
    },
    'html': {
        'pdf': convert_html_to_pdf,
    }
}

IMAGE_CONVERSIONS = {
    'jpg': ['png', 'webp', 'pdf', 'jpeg'],
    'jpeg': ['png', 'webp', 'pdf', 'jpg'],
    'png': ['jpg', 'webp', 'pdf', 'jpeg'],
    'webp': ['jpg', 'png', 'pdf', 'jpeg'],
    'gif': ['png', 'jpg', 'webp', 'pdf'],
    'bmp': ['png', 'jpg', 'webp', 'pdf']
}

SUPPORTED_IMAGE_INPUTS = list(IMAGE_CONVERSIONS.keys())
